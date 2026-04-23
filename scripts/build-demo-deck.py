#!/usr/bin/env python3
"""Build the Netchex Visitor Signals demo deck.

Produces docs/demo-deck.pptx. Screenshots are pulled from public/screenshots/
if present; otherwise the slide renders with a placeholder box so the deck
still opens cleanly before you've captured anything.

Run:
    python3 scripts/build-demo-deck.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent.parent
SHOTS = ROOT / "public" / "screenshots"
OUT = ROOT / "docs" / "demo-deck.pptx"

# Zinc palette — black/white/zinc only, per UI guidelines.
BLACK = RGBColor(0x09, 0x09, 0x0B)       # zinc-950
ZINC_800 = RGBColor(0x27, 0x27, 0x2A)
ZINC_600 = RGBColor(0x52, 0x52, 0x5B)
ZINC_400 = RGBColor(0xA1, 0xA1, 0xAA)
ZINC_200 = RGBColor(0xE4, 0xE4, 0xE7)
ZINC_100 = RGBColor(0xF4, 0xF4, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_background(slide, color=WHITE):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    # send to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_text(slide, left, top, width, height, text, *, size=18, bold=False,
             color=BLACK, align=PP_ALIGN.LEFT, font="Inter"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for line in text.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_accent_bar(slide, x=Inches(0.6), y=Inches(0.75), h=Inches(0.08), w=Inches(0.6)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLACK
    return bar


def add_image_or_placeholder(slide, name, left, top, width, height, label):
    path = SHOTS / f"{name}.png"
    if path.exists():
        slide.shapes.add_picture(str(path), left, top, width=width, height=height)
        return
    # placeholder
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = ZINC_100
    box.line.color.rgb = ZINC_200
    box.line.width = Pt(1)
    tb = box.text_frame
    tb.margin_left = tb.margin_right = Inches(0.2)
    tb.margin_top = tb.margin_bottom = Inches(0.2)
    tb.word_wrap = True
    p = tb.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = f"[ screenshot: {name}.png ]\n{label}"
    r.font.name = "Inter"
    r.font.size = Pt(14)
    r.font.color.rgb = ZINC_400
    r.font.italic = True


def slide_title_subtitle(prs, title, subtitle, footer):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)

    # black block on left
    block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(5.0), SLIDE_H)
    block.line.fill.background()
    block.fill.solid()
    block.fill.fore_color.rgb = BLACK

    add_text(slide, Inches(0.6), Inches(0.7), Inches(4), Inches(0.4),
             "NETCHEX  /  VISITOR SIGNALS", size=11, bold=True, color=ZINC_400)

    add_text(slide, Inches(0.6), Inches(2.5), Inches(4.2), Inches(2.6),
             title, size=44, bold=True, color=WHITE)
    add_text(slide, Inches(0.6), Inches(5.0), Inches(4.2), Inches(1.4),
             subtitle, size=16, color=ZINC_200)

    add_text(slide, Inches(5.6), Inches(0.9), Inches(7.2), Inches(0.4),
             "Demo walkthrough", size=11, bold=True, color=ZINC_600)
    add_accent_bar(slide, Inches(5.6), Inches(1.35), Inches(0.06), Inches(0.8))

    add_text(slide, Inches(5.6), Inches(2.2), Inches(7.2), Inches(3.4),
             "Anonymous website traffic → identified visitor\n"
             "→ ICP-scored lead → personalized outbound email\n"
             "→ simulated Gong Engage send.\n\n"
             "One pipeline. Five integrations. Zero black boxes.",
             size=20, color=ZINC_800)

    add_text(slide, Inches(5.6), Inches(6.8), Inches(7.2), Inches(0.4),
             footer, size=10, color=ZINC_400)
    return slide


def slide_standard(prs, section, title, body_lines, screenshot=None,
                    screenshot_label=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)

    add_text(slide, Inches(0.6), Inches(0.45), Inches(6), Inches(0.35),
             section.upper(), size=10, bold=True, color=ZINC_400)
    add_accent_bar(slide)
    add_text(slide, Inches(0.6), Inches(1.05), Inches(12), Inches(0.9),
             title, size=32, bold=True, color=BLACK)

    # bullet column
    body_left = Inches(0.6)
    body_top = Inches(2.2)
    body_w = Inches(5.0) if screenshot else Inches(11.0)
    body_h = Inches(5.0)

    tb = slide.shapes.add_textbox(body_left, body_top, body_w, body_h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    first = True
    for line in body_lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(10)
        r = p.add_run()
        if line.startswith("**"):
            # bold header line
            text = line.strip("*")
            r.text = text
            r.font.bold = True
            r.font.size = Pt(17)
            r.font.color.rgb = BLACK
        else:
            r.text = f"•  {line}"
            r.font.size = Pt(15)
            r.font.color.rgb = ZINC_800
        r.font.name = "Inter"

    if screenshot:
        add_image_or_placeholder(
            slide,
            screenshot,
            Inches(6.0),
            Inches(2.2),
            Inches(6.8),
            Inches(4.6),
            screenshot_label,
        )

    add_text(slide, Inches(0.6), Inches(7.1), Inches(12), Inches(0.3),
             "Netchex Visitor Signals  —  demo deck", size=9, color=ZINC_400)
    return slide


def slide_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BLACK)
    add_text(slide, Inches(0.6), Inches(0.7), Inches(12), Inches(0.4),
             "WHAT'S NEXT", size=11, bold=True, color=ZINC_400)

    add_text(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(1.2),
             "From demo to production", size=40, bold=True, color=WHITE)

    roadmap = [
        ("Real Gong Engage send", "Behind a feature flag. Awaiting Legal sign-off."),
        ("Auto-assign to AE", "Gong's CRM mapping already tells us the owner."),
        ("Slack alert for Tier 1 leads", "Same hook, different endpoint."),
        ("Daily SDR digest", "Top 10 qualified visitors from the last 24h."),
        ("Monitoring dashboard", "Time-to-qualification, LLM spend, enrichment credits."),
    ]

    y = Inches(3.0)
    for title, desc in roadmap:
        add_text(slide, Inches(0.6), y, Inches(4.5), Inches(0.4),
                 title, size=16, bold=True, color=WHITE)
        add_text(slide, Inches(5.2), y, Inches(7.6), Inches(0.4),
                 desc, size=14, color=ZINC_400)
        y += Inches(0.65)

    add_text(slide, Inches(0.6), Inches(6.9), Inches(12), Inches(0.3),
             "Questions?", size=14, bold=True, color=WHITE)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Slide 1 — title
    slide_title_subtitle(
        prs,
        "Turn anonymous traffic\ninto vetted outreach",
        "A single pipeline from RB2B pixel hit to a personalized email "
        "sitting in front of an SDR — in under 15 seconds and under 15 cents.",
        "Presented by the Sales Ops team",
    )

    # Slide 2 — the problem
    slide_standard(
        prs,
        section="The gap",
        title="97% of visitors never fill out a form",
        body_lines=[
            "**Identified ≠ qualified.**",
            "RB2B surfaces ~40% of visitors by name, email, company, LinkedIn.",
            "Most of those aren't buyers — wrong title, wrong industry, wrong size.",
            "An SDR reading the feed manually burns an hour a day and still misses the real leads.",
            "**The bottleneck is triage and personalization — not identification.**",
        ],
    )

    # Slide 3 — pipeline
    slide_standard(
        prs,
        section="Architecture",
        title="Five steps, fully configurable",
        body_lines=[
            "**RB2B  →  Clay  →  ICP score  →  SF/Gong dedup  →  LLM draft  →  SDR**",
            "RB2B webhook (or CSV bulk) drops identified visitors into Postgres.",
            "Clay enriches on LinkedIn URL → work email, title, company firmographics.",
            "ICP engine scores 0–100 against editable rules; tiers leads 1/2/3.",
            "Before drafting, we check Salesforce + Gong for existing deals or prior contact.",
            "Gemini 2.5 Flash drafts a voice-matched email; SDR copies or simulates send.",
        ],
        screenshot="02-dashboard",
        screenshot_label="Dashboard — stats + qualified visitor table",
    )

    # Slide 4 — ICP scoring
    slide_standard(
        prs,
        section="Scoring the funnel",
        title="The ICP lives in the app, not in code",
        body_lines=[
            "**Four signals, 100 points total.**",
            "Firmographics (40) — industry, size 50–500, US-based.",
            "Intent (30) — which product pages; /pricing and /demo weigh extra.",
            "Persona fit (20) — HR / payroll / ops / owner titles.",
            "Capacity (10) — revenue band, repeat visits.",
            "**Sales Ops edits the rules in Settings.** One-click re-score on the top 50.",
        ],
        screenshot="03-score",
        screenshot_label="Visitor detail — Score Reasoning popover",
    )

    # Slide 5 — CRM dedup
    slide_standard(
        prs,
        section="Don't step on a live deal",
        title="Every email is checked against Salesforce + Gong first",
        body_lines=[
            "**One click: Check Salesforce & Gong.**",
            "SOQL lookup by email — Lead or Contact, with owner + last activity.",
            "Gong API — prior call count, prior email count, last touchpoint.",
            "Results cache on the visitor row; repeat dashboard loads are free.",
            "If this person is already in-play, the SDR gets a warning — not a draft.",
        ],
        screenshot="04-crm",
        screenshot_label="Visitor detail — CRM & Prior Touchpoints card",
    )

    # Slide 6 — email generation
    slide_standard(
        prs,
        section="The tone is the product",
        title="Voice-matched email, regenerable on demand",
        body_lines=[
            "**Gemini 2.5 Flash, ~3s, ~$0.0005 per draft.**",
            "Prompt reverse-engineered from 10 real SDR emails (Teresa, Scott).",
            "Greeting: \"Hey [first name],\" — never \"Dear\".",
            "A real Netchex customer name (IVY, CFAN, Willbanks) when the industry matches.",
            "CTA always ends with a duration: \"Is a 10-minute overview worth scheduling?\"",
            "Marketing retunes the voice by editing the template — no deploy.",
        ],
        screenshot="05-outreach",
        screenshot_label="Generated outreach email + Copy / Send controls",
    )

    # Slide 7 — closing
    slide_closing(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"wrote {OUT.relative_to(ROOT)}")


if __name__ == "__main__":
    build()
